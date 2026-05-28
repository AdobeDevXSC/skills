"""
Build Bank of America EDS POV presentation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ────────────────────────────────────────────────────────────
RED      = RGBColor(0xFA, 0x0F, 0x00)   # Adobe red
DARK     = RGBColor(0x1C, 0x1C, 0x1C)   # Near-black
MID      = RGBColor(0x46, 0x46, 0x46)   # Dark-grey body
LIGHT    = RGBColor(0xF5, 0xF5, 0xF5)   # Off-white bg
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BLUE     = RGBColor(0x14, 0x73, 0xE6)   # Accent blue
AMBER    = RGBColor(0xE6, 0x8A, 0x14)   # Warm accent
GREEN    = RGBColor(0x26, 0x8E, 0x6C)   # Positive green
SLATE    = RGBColor(0x4A, 0x5C, 0x6A)   # Table header slate
ROW_ALT  = RGBColor(0xED, 0xF2, 0xF7)   # Alternating row tint
ROW_INT  = RGBColor(0xFF, 0xF3, 0xCD)   # Internal-only highlight (amber tint)

# ── Slide dimensions (16:9 widescreen) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Fonts ────────────────────────────────────────────────────────────────────
FONT = "Calibri"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # blank layout

# ─────────────────────────────────────────────────────────────────────────────
# Helper utilities
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font=FONT, size=12, bold=False, italic=False,
                color=DARK, align=PP_ALIGN.LEFT,
                wrap=True, valign=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, size=11, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None):
    """Red top bar with white title + optional subtitle."""
    add_rect(slide, 0, 0, W, Inches(1.05), fill=RED)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.12), Inches(11.5), Inches(0.65),
                size=22, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.68), Inches(11.5), Inches(0.34),
                    size=12, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC))
    # bottom rule
    add_rect(slide, 0, H - Inches(0.28), W, Inches(0.05), fill=RED)
    # footer
    add_textbox(slide,
                "Adobe Confidential  |  Bank of America POV  |  May 2026",
                Inches(0.3), H - Inches(0.26), Inches(9), Inches(0.24),
                size=8, color=MID)
    add_textbox(slide, "adobe.com",
                W - Inches(1.4), H - Inches(0.26), Inches(1.1), Inches(0.24),
                size=8, color=RED, align=PP_ALIGN.RIGHT)


def make_table(slide, headers, rows,
               x, y, w, h,
               header_fill=SLATE, alt_fill=ROW_ALT,
               header_color=WHITE, body_color=DARK,
               col_widths=None, header_size=10, body_size=9,
               highlight_rows=None):
    """Draw a styled table without using pptx native tables (more layout control)."""
    ncols = len(headers)
    nrows = len(rows)
    total_rows = nrows + 1   # +1 for header

    if col_widths is None:
        cw = [w / ncols] * ncols
    else:
        cw = [Inches(c) for c in col_widths]

    row_h = h / total_rows

    # Header row
    cx = x
    for ci, hdr in enumerate(headers):
        add_rect(slide, cx, y, cw[ci], row_h, fill=header_fill)
        add_textbox(slide, hdr,
                    cx + Inches(0.06), y + Inches(0.04),
                    cw[ci] - Inches(0.12), row_h - Inches(0.06),
                    size=header_size, bold=True, color=header_color)
        cx += cw[ci]

    # Data rows
    for ri, row in enumerate(rows):
        ry = y + row_h * (ri + 1)
        bg = ROW_INT if (highlight_rows and ri in highlight_rows) else \
             (alt_fill if ri % 2 == 0 else WHITE)
        cx = x
        for ci, cell in enumerate(row):
            add_rect(slide, cx, ry, cw[ci], row_h, fill=bg,
                     line=RGBColor(0xD0, 0xD0, 0xD0))
            add_textbox(slide, str(cell),
                        cx + Inches(0.06), ry + Inches(0.03),
                        cw[ci] - Inches(0.12), row_h - Inches(0.06),
                        size=body_size, color=body_color, wrap=True)
            cx += cw[ci]


def score_badge(slide, label, score, x, y, color=BLUE):
    """A circular-ish score badge."""
    add_rect(slide, x, y, Inches(1.2), Inches(1.2), fill=color)
    add_textbox(slide, str(score),
                x, y + Inches(0.15), Inches(1.2), Inches(0.55),
                size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
                x, y + Inches(0.72), Inches(1.2), Inches(0.38),
                size=8, color=WHITE, align=PP_ALIGN.CENTER)


def kpi_box(slide, label, value, sub, x, y, w=Inches(2.4), h=Inches(1.1), color=BLUE):
    add_rect(slide, x, y, w, h, fill=color)
    add_textbox(slide, value,
                x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.5),
                size=22, bold=True, color=WHITE)
    add_textbox(slide, label,
                x + Inches(0.1), y + Inches(0.56), w - Inches(0.2), Inches(0.28),
                size=9, bold=True, color=WHITE)
    add_textbox(slide, sub,
                x + Inches(0.1), y + Inches(0.78), w - Inches(0.2), Inches(0.26),
                size=8, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF))


def bullet_box(slide, title, bullets, x, y, w, h,
               title_color=RED, bullet_size=10, title_size=12):
    add_textbox(slide, title, x, y, w, Inches(0.3),
                size=title_size, bold=True, color=title_color)
    txb = slide.shapes.add_textbox(x, y + Inches(0.32), w, h - Inches(0.32))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        add_para(tf, f"• {b}", size=bullet_size, color=MID,
                 space_before=(2 if i > 0 else 0))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)

# Full dark background
add_rect(slide, 0, 0, W, H, fill=DARK)

# Red accent bar left
add_rect(slide, 0, 0, Inches(0.18), H, fill=RED)

# Red stripe near bottom
add_rect(slide, 0, H - Inches(1.4), W, Inches(0.06), fill=RED)

# Adobe wordmark area (simulated)
add_textbox(slide, "adobe",
            Inches(0.4), Inches(0.35), Inches(2.5), Inches(0.5),
            size=28, bold=True, color=RED)

# Eyebrow
add_textbox(slide, "SOLUTION CONSULTANT  •  POINT OF VIEW",
            Inches(0.4), Inches(1.3), Inches(11), Inches(0.35),
            size=10, color=RGBColor(0xAA, 0xAA, 0xAA), bold=False)

# Main title
add_textbox(slide, "Bank of America",
            Inches(0.4), Inches(1.75), Inches(10), Inches(1.1),
            size=48, bold=True, color=WHITE)

# Subtitle
add_textbox(slide,
            "AEM Edge Delivery Services + Document Authoring",
            Inches(0.4), Inches(2.85), Inches(10), Inches(0.55),
            size=22, bold=False, color=RGBColor(0xCC, 0xCC, 0xCC))

# Divider line
add_rect(slide, Inches(0.4), Inches(3.55), Inches(5), Inches(0.03), fill=RED)

# Meta line
add_textbox(slide, "Prepared: May 27, 2026  |  Adobe Solution Consulting",
            Inches(0.4), Inches(3.72), Inches(9), Inches(0.3),
            size=11, color=RGBColor(0x99, 0x99, 0x99))

# Tagline
add_textbox(slide,
            "\"Digital is the centerpiece of our relationship-driven strategy.\"\n— Nikki Katz, Head of Digital, Bank of America",
            Inches(0.4), Inches(4.3), Inches(9), Inches(0.7),
            size=12, italic=True, color=RGBColor(0xFF, 0xCC, 0xCC))

# Bottom strip text
add_textbox(slide, "Adobe Confidential",
            Inches(0.4), H - Inches(1.25), Inches(5), Inches(0.3),
            size=9, color=RGBColor(0x88, 0x88, 0x88))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — EXECUTIVE SUMMARY
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Executive Summary", "The case for EDS at Bank of America")

# 3 columns
col_w = Inches(3.9)
col_h = Inches(4.6)
col_y = Inches(1.25)
gap   = Inches(0.25)

cols = [
    {
        "icon": "🏛",
        "color": BLUE,
        "title": "The Situation",
        "body": [
            "BofA runs one of the largest AEM estates in the world — 16+ DRIDs, $23M+ Sites ARR.",
            "bankofamerica.com is a complex SPA-based banking portal. Its EDS Migration Score is 32 — the portal itself is NOT the target.",
            "But the editorial layer — newsroom, Better Money Habits, and about.bankofamerica.com — is a natural EDS fit.",
        ]
    },
    {
        "icon": "✅",
        "color": GREEN,
        "title": "The Proof Point",
        "body": [
            "BofA has already validated EDS internally.",
            "Per FLM notes: \"Neushler & team rebuilt some of their sites on EDS.\"",
            "Phase 0 is done. The conversation is not 'should we try EDS?' — it is 'how do we formalize what's already working and scale it?'",
        ]
    },
    {
        "icon": "💰",
        "color": AMBER,
        "title": "The Opportunity",
        "body": [
            "$9M GNARR potential.",
            "Renewal: December 31, 2027 — 18 months to build a compelling success story on the editorial layer.",
            "Three editorial properties + campaign microsites = ~1,000–1,500 pages ready for EDS + Document Authoring.",
        ]
    },
]

for i, col in enumerate(cols):
    cx = Inches(0.3) + i * (col_w + gap)
    add_rect(slide, cx, col_y, col_w, col_h, fill=WHITE,
             line=RGBColor(0xD8, 0xD8, 0xD8))
    # color bar top
    add_rect(slide, cx, col_y, col_w, Inches(0.08), fill=col["color"])
    # title
    add_textbox(slide, col["title"],
                cx + Inches(0.18), col_y + Inches(0.2), col_w - Inches(0.36), Inches(0.4),
                size=14, bold=True, color=col["color"])
    # bullets
    txb = slide.shapes.add_textbox(
        cx + Inches(0.18), col_y + Inches(0.7),
        col_w - Inches(0.36), col_h - Inches(0.85))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for j, b in enumerate(col["body"]):
        add_para(tf, f"• {b}", size=10.5, color=MID,
                 space_before=(6 if j > 0 else 0))

# Bottom callout bar
add_rect(slide, 0, Inches(6.1), W, Inches(0.75), fill=RED)
add_textbox(slide,
            "The editorial layer is a low-risk, high-visibility EDS beachhead — and BofA's own team has already proven it works.",
            Inches(0.4), Inches(6.16), Inches(12.5), Inches(0.55),
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — CUSTOMER SNAPSHOT
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Customer Snapshot", "Bank of America  •  Financial Services")

# KPI row
kpis = [
    ("Current Sites ARR", "$3.2M", "On-prem portion",          BLUE),
    ("Full Adobe Relationship", "$23M+ Sites", "$29M+ total",  SLATE),
    ("Open Pipeline", "$9M",    "GNARR potential",             AMBER),
    ("Renewal Date",  "Dec 2027", "FY28 Q1 — 18 months",       GREEN),
    ("Priority Score", "0",     "Top priority (active)",       RED),
]
kx = Inches(0.3)
for label, value, sub, color in kpis:
    kpi_box(slide, label, value, sub, kx, Inches(1.2), Inches(2.5), Inches(1.1), color)
    kx += Inches(2.55)

# Scores section
add_textbox(slide, "Readiness Scores",
            Inches(0.3), Inches(2.55), Inches(5), Inches(0.3),
            size=11, bold=True, color=DARK)

score_badge(slide, "EDS Migration\nScore", 32,  Inches(0.3),  Inches(2.9), color=AMBER)
score_badge(slide, "LLM Visibility\nScore", 80, Inches(1.65), Inches(2.9), color=GREEN)

add_textbox(slide,
            "Score: 32 / 100  ⚠️  Significant headwinds on the core banking portal",
            Inches(3.0), Inches(2.95), Inches(5.5), Inches(0.3),
            size=9.5, bold=True, color=AMBER)
add_textbox(slide,
            "\"Largest US consumer bank sharing the same complex AEM environment as JPMorgan\n"
            "(16 DRIDs); consumer banking portals and financial tools dominate, leaving limited\n"
            "marketing-only EDS scope.\"  →  Pitch targets the editorial layer, not the portal.",
            Inches(3.0), Inches(3.28), Inches(5.5), Inches(0.65),
            size=9, italic=True, color=MID)

add_textbox(slide,
            "Score: 80 / 100  ✅  Strong AI search visibility",
            Inches(3.0), Inches(4.02), Inches(5.5), Inches(0.3),
            size=9.5, bold=True, color=GREEN)
add_textbox(slide,
            "BofA captures ~22 AI platform mentions, outpacing First Citizens and competing\n"
            "directly with Chase/JPMorgan in AI-driven banking discovery.",
            Inches(3.0), Inches(4.32), Inches(5.5), Inches(0.4),
            size=9, italic=True, color=MID)

# Right panel — table
tbl_headers = ["Field", "Detail"]
tbl_rows = [
    ["Industry",          "Financial Services (FSI)"],
    ["CMS / Platform",    "AEM On-Premises  (16+ DRIDs, SPA architecture)"],
    ["Products in Use",   "AEM Sites + Assets (On-Prem), Analytics→CJA, Target, AEP"],
    ["M2EDS Status",      "In Progress — \"Neushler & team rebuilt sites on EDS\""],
    ["Pipeline Wave",     "In-Progress  (FLM: Herzenberg)"],
    ["Region",            "FSI (Financial Services Industry)"],
    ["Web Presence",      "bankofamerica.com (27 sitemap sections), newsroom, BMH, about"],
]
make_table(slide, tbl_headers, tbl_rows,
           Inches(8.7), Inches(1.2),
           Inches(4.4), Inches(5.2),
           col_widths=[1.5, 2.8],
           header_size=9, body_size=8.5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — BUSINESS CONTEXT
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Business Context & Strategic Priorities",
             "Bank of America  •  $13B annual technology investment")

priorities = [
    {
        "num": "01",
        "title": "AI-First Experience Delivery",
        "color": RED,
        "body": (
            "$13B/year tech spend  |  $4B on new initiatives  |  1,400 AI patents\n"
            "Erica AI assistant: 700M+ interactions from 20.6M users\n"
            "17,000 programmers using AI coding tools  |  250+ AI models in production\n"
            "\"Digital is the centerpiece of our relationship-driven strategy.\"\n"
            "— Nikki Katz, Head of Digital"
        ),
    },
    {
        "num": "02",
        "title": "Client Engagement at Scale",
        "color": BLUE,
        "body": (
            "59M verified digital users  |  30B interactions in 2025 (+14% YoY)\n"
            "16.6B logins  |  13.3B alerts sent\n"
            "#1 J.D. Power: Mobile Banking, Mortgage Servicer Digital, Online Banking\n"
            "Publishing infrastructure must keep pace — current AEM on-prem creates bottlenecks."
        ),
    },
    {
        "num": "03",
        "title": "Financial Education as a Growth Channel",
        "color": GREEN,
        "body": (
            "Better Money Habits (bettermoneyhabits.bankofamerica.com): 300+ articles,\n"
            "9 topic areas, bilingual (EN/ES), active since 2013.\n"
            "Primary AI-indexed, search-discoverable content surface.\n"
            "Key acquisition + loyalty lever — today blocked by AEM publishing overhead."
        ),
    },
]

pw = Inches(4.1)
ph = Inches(4.5)
px = Inches(0.22)

for i, p in enumerate(priorities):
    cx = px + i * (pw + Inches(0.15))
    add_rect(slide, cx, Inches(1.2), pw, ph, fill=WHITE,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(slide, cx, Inches(1.2), pw, Inches(0.06), fill=p["color"])
    # number badge
    add_rect(slide, cx + Inches(0.15), Inches(1.28), Inches(0.55), Inches(0.55),
             fill=p["color"])
    add_textbox(slide, p["num"],
                cx + Inches(0.15), Inches(1.3), Inches(0.55), Inches(0.5),
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # title
    add_textbox(slide, p["title"],
                cx + Inches(0.82), Inches(1.33), pw - Inches(1.0), Inches(0.45),
                size=12, bold=True, color=p["color"])
    # divider
    add_rect(slide, cx + Inches(0.15), Inches(1.88), pw - Inches(0.3),
             Inches(0.02), fill=RGBColor(0xDD, 0xDD, 0xDD))
    # body
    add_textbox(slide, p["body"],
                cx + Inches(0.18), Inches(1.96), pw - Inches(0.36), ph - Inches(0.85),
                size=10, color=MID, wrap=True)

# Source bar
add_rect(slide, 0, Inches(5.88), W, Inches(0.28), fill=RGBColor(0xEE, 0xEE, 0xEE))
add_textbox(slide,
            "Sources: BofA Digital Innovations Press Release (March 2026)  |  "
            "BofA Innovation Awards (Oct 2025)  |  BofA Better Money Habits Fact Sheet 2024",
            Inches(0.3), Inches(5.91), Inches(12), Inches(0.24),
            size=7.5, italic=True, color=MID)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — CURRENT STATE & PAIN POINTS
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Current State Assessment", "Pain points on the editorial layer")

# Left: site properties
add_textbox(slide, "Editorial Properties (EDS Target Layer)",
            Inches(0.3), Inches(1.2), Inches(5.6), Inches(0.3),
            size=11, bold=True, color=DARK)

props = [
    ("newsroom.\nbankofamerica.com",  BLUE,  "Daily press releases\nAEM/JCR-based\n209+ releases, exec bios,\nresearch reports"),
    ("bettermoneyhabits.\nbankofamerica.com", GREEN, "300+ articles, 9 topics\nEN + ES (bilingual)\nInteractive tools\nAEM/JCR-based"),
    ("about.\nbankofamerica.com",     SLATE, "Sustainability, investor\nresources, fact sheets\n~150+ pages\nAEM-based"),
]

for i, (name, color, desc) in enumerate(props):
    bx = Inches(0.3) + i * Inches(1.9)
    add_rect(slide, bx, Inches(1.6), Inches(1.75), Inches(2.2), fill=WHITE,
             line=color)
    add_rect(slide, bx, Inches(1.6), Inches(1.75), Inches(0.06), fill=color)
    add_textbox(slide, name, bx + Inches(0.1), Inches(1.72),
                Inches(1.55), Inches(0.55),
                size=8.5, bold=True, color=color, wrap=True)
    add_textbox(slide, desc, bx + Inches(0.1), Inches(2.3),
                Inches(1.55), Inches(1.4),
                size=8, color=MID, wrap=True)

# Right: pain points table
add_textbox(slide, "Identified Pain Points",
            Inches(6.0), Inches(1.2), Inches(7.0), Inches(0.3),
            size=11, bold=True, color=DARK)

pt_headers = ["Pain Point", "Evidence", "Business Impact"]
pt_rows = [
    ["Editorial teams blocked\non banking platform",
     "newsroom + BMH run on AEM on-prem\nbuilt for banking portals, not publishing",
     "Publishing delays; IT\ndependency for minor updates"],
    ["Content velocity\nbottleneck",
     "11 Target activities launched in 2026;\nactive digital program needs fast publish",
     "Slow campaign execution;\ndev-gated content updates"],
    ["AEM on-prem lifecycle\nrisk",
     "16 DRIDs on-prem; Dec 2027 renewal;\nAEM CS is the industry trajectory",
     "Renewal leverage; competitors\ngain on Core Web Vitals + AI search"],
    ["SPA architecture\nperformance ceiling",
     "spa-assets URL pattern; SPA content\nless accessible to AI search engines",
     "LLM Score 80 — editorial layer\ncould score higher with EDS HTML"],
    ["Fragmented authoring\nacross properties",
     "newsroom, BMH, about are separate\nAEM deployments with separate pipelines",
     "High operational overhead;\nduplicated infrastructure cost"],
]
make_table(slide, pt_headers, pt_rows,
           Inches(6.0), Inches(1.6),
           Inches(7.1), Inches(4.5),
           col_widths=[1.8, 2.5, 2.7],
           header_size=9, body_size=8.5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — MIGRATION SCOPE ESTIMATE
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Migration Scope Estimate",
             "Editorial layer only — core banking portal excluded")

# Page count boxes
add_textbox(slide, "Page Count  (Phase 1 Editorial Scope)",
            Inches(0.3), Inches(1.2), Inches(8), Inches(0.3),
            size=11, bold=True, color=DARK)

pc_items = [
    ("newsroom.\nbankofamerica.com",  "500+",  "Press releases,\nexec bios, research",   BLUE),
    ("bettermoneyhabits.\nbankofamerica.com", "300+",  "Articles, tools,\nbilingual",           GREEN),
    ("about.\nbankofamerica.com",      "150+",  "Sustainability, investor,\nfact sheets",       SLATE),
    ("Campaign\nmicrosites",           "50+/yr", "Ongoing net-new\ncampaign sites",            AMBER),
    ("Total EDS\nScope",               "~1,000–\n1,500",  "Phase 1\nestimate",                 RED),
]

for i, (label, count, sub, color) in enumerate(pc_items):
    bx = Inches(0.28) + i * Inches(2.55)
    add_rect(slide, bx, Inches(1.6), Inches(2.35), Inches(1.5), fill=WHITE,
             line=color)
    add_rect(slide, bx, Inches(1.6), Inches(2.35), Inches(0.05), fill=color)
    add_textbox(slide, label, bx + Inches(0.1), Inches(1.7),
                Inches(2.1), Inches(0.45),
                size=8, bold=True, color=color, wrap=True)
    add_textbox(slide, count, bx + Inches(0.1), Inches(2.1),
                Inches(2.1), Inches(0.52),
                size=20, bold=True, color=DARK)
    add_textbox(slide, sub, bx + Inches(0.1), Inches(2.6),
                Inches(2.1), Inches(0.38),
                size=7.5, color=MID, wrap=True)

# Block inventory
add_textbox(slide, "Block Inventory  (sampled from newsroom + Better Money Habits)",
            Inches(0.3), Inches(3.3), Inches(8), Inches(0.3),
            size=11, bold=True, color=DARK)

bi_headers = ["Block Type", "Found On", "EDS Block Collection", "Status"]
bi_rows = [
    ["Hero / banner",            "All editorial homepages",         "✅ Yes",  "Adopt"],
    ["Rich text / article body", "Press releases, BMH articles",    "✅ Yes",  "Adopt"],
    ["Card grid / article list", "Category hubs, topic pages",      "✅ Yes",  "Adopt"],
    ["Navigation header/footer", "All pages",                       "✅ Yes",  "Adopt"],
    ["Video embed",              "BMH video articles, newsroom",     "✅ Yes",  "Adopt"],
    ["Related content carousel", "BMH detail, newsroom",            "✅ Yes",  "Adopt"],
    ["Social share / subscription", "Newsroom, BMH",               "⚡ Partial", "Customize"],
    ["Language toggle (EN/ES)",  "BMH localization",                "⚡ Partial", "Customize"],
    ["Search + filter",          "Newsroom search, BMH topic search", "❌ No", "Custom build"],
    ["Interactive financial tool","BMH calculators, quiz",           "❌ No", "Custom build"],
]
make_table(slide, bi_headers, bi_rows,
           Inches(0.3), Inches(3.7),
           Inches(12.7), Inches(2.9),
           col_widths=[2.4, 3.0, 2.8, 1.85],
           header_size=9, body_size=8.5,
           highlight_rows=[8, 9])

# Complexity rating
add_rect(slide, Inches(9.8), Inches(1.2), Inches(3.2), Inches(0.9), fill=GREEN)
add_textbox(slide, "Migration Complexity",
            Inches(9.85), Inches(1.22), Inches(3.1), Inches(0.3),
            size=9, bold=True, color=WHITE)
add_textbox(slide, "LOW – MEDIUM",
            Inches(9.85), Inches(1.5), Inches(3.1), Inches(0.5),
            size=20, bold=True, color=WHITE)

add_textbox(slide, "20% custom block ratio  |  2 custom blocks out of 10 total",
            Inches(9.8), Inches(2.2), Inches(3.2), Inches(0.35),
            size=9, bold=True, color=GREEN)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — THE EDS OPPORTUNITY
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "The Opportunity: AEM Edge Delivery Services + Document Authoring",
             "Why EDS — and why now")

# Two-column layout
col_w = Inches(6.3)

# LEFT: Why EDS
add_rect(slide, Inches(0.25), Inches(1.2), col_w, Inches(5.4), fill=WHITE,
         line=RGBColor(0xDD, 0xDD, 0xDD))
add_rect(slide, Inches(0.25), Inches(1.2), col_w, Inches(0.07), fill=BLUE)
add_textbox(slide, "Why EDS Is the Right Answer for Bank of America",
            Inches(0.45), Inches(1.33), col_w - Inches(0.3), Inches(0.38),
            size=12, bold=True, color=BLUE)

why_eds = [
    ("Lighthouse 100 for the editorial layer",
     "bankofamerica.com's SPA architecture has a performance ceiling. EDS delivers "
     "Lighthouse 100 natively for content pages — newsroom and BMH would see dramatic "
     "Core Web Vitals improvements."),
    ("AI search visibility gain on Better Money Habits",
     "EDS produces semantically clean, static HTML — the format AI engines (ChatGPT, "
     "Perplexity, Google AI Overviews) parse best. Moving 300+ BMH financial education "
     "articles to EDS strengthens BofA's AI search authority in personal finance."),
    ("Phase 0 is done — Neushler's team proved it",
     "Internal EDS sites already exist at BofA. This is not a leap of faith; "
     "it is a formalization of work that has already succeeded."),
    ("Renewal readiness",
     "Dec 2027 renewal = 18-month window. EDS Phase 1 (newsroom) + Phase 2 (BMH) "
     "gives the account team a measured success story to anchor the $9M GNARR renewal."),
]

yy = Inches(1.82)
for title, body in why_eds:
    add_textbox(slide, f"▸  {title}",
                Inches(0.45), yy, col_w - Inches(0.3), Inches(0.26),
                size=10, bold=True, color=DARK)
    add_textbox(slide, body,
                Inches(0.65), yy + Inches(0.26), col_w - Inches(0.5), Inches(0.42),
                size=9, color=MID, wrap=True)
    yy += Inches(0.78)

# RIGHT: Why Document Authoring
rx = Inches(6.78)
add_rect(slide, rx, Inches(1.2), col_w, Inches(5.4), fill=WHITE,
         line=RGBColor(0xDD, 0xDD, 0xDD))
add_rect(slide, rx, Inches(1.2), col_w, Inches(0.07), fill=RED)
add_textbox(slide, "Why Document Authoring on da.live Is Specifically Right",
            rx + Inches(0.2), Inches(1.33), col_w - Inches(0.3), Inches(0.38),
            size=12, bold=True, color=RED)

why_da = [
    ("Authors already live in Microsoft Word + SharePoint",
     "BofA runs Microsoft 365 at enterprise scale. Document Authoring meets the "
     "newsroom and BMH editors where they already work — no new tools to learn."),
    ("Press release in < 30 minutes",
     "Today: draft → AEM ticket → component editing → deployment → publish (days). "
     "With Document Authoring: draft in Word → save to SharePoint → Sidekick click → live (< 30 min)."),
    ("Bilingual publishing simplified",
     "BMH's EN/ES content today requires engineering involvement. With Document "
     "Authoring, separate SharePoint folders per locale handle everything. No developer needed."),
    ("SharePoint governance = financial services compliant",
     "da.live uses SharePoint permission models — the same access control + audit "
     "trail infrastructure BofA already uses for enterprise documents."),
]

yy = Inches(1.82)
for title, body in why_da:
    add_textbox(slide, f"▸  {title}",
                rx + Inches(0.2), yy, col_w - Inches(0.3), Inches(0.26),
                size=10, bold=True, color=DARK)
    add_textbox(slide, body,
                rx + Inches(0.4), yy + Inches(0.26), col_w - Inches(0.5), Inches(0.42),
                size=9, color=MID, wrap=True)
    yy += Inches(0.78)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — MIGRATION APPROACH
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Recommended Migration Approach", "Phase-by-phase — editorial layer first")

phases = [
    {
        "phase": "Phase 0",
        "label": "Validate & Inventory",
        "time": "Now  (2–4 wks)",
        "color": SLATE,
        "scope": "Identify & document the sites Neushler's team already built on EDS. Run Lighthouse. Declare Phase 0 complete.",
        "outcome": "Internal case study. EDS already live at BofA. Lighthouse scores captured.",
    },
    {
        "phase": "Phase 1",
        "label": "Newsroom on Document Authoring",
        "time": "6–8 weeks",
        "color": BLUE,
        "scope": "Reproduce newsroom.bankofamerica.com on EDS + Document Authoring. 6 standard blocks. Communications team authors from Word/SharePoint.",
        "outcome": "Lighthouse 100. Press releases publish from Word in < 30 min. Champion in communications leadership.",
    },
    {
        "phase": "Phase 2",
        "label": "Better Money Habits Migration",
        "time": "8–12 weeks",
        "color": GREEN,
        "scope": "Migrate bettermoneyhabits.bankofamerica.com. 300+ articles, bilingual, 2 custom blocks (search + interactive tools).",
        "outcome": "Full financial education platform on EDS. AI-indexed structured HTML. EN/ES publishing via SharePoint.",
    },
    {
        "phase": "Phase 3",
        "label": "About + Campaign Microsites",
        "time": "Ongoing",
        "color": AMBER,
        "scope": "Migrate about.bankofamerica.com. Establish EDS as the default platform for all net-new campaign + product launch sites.",
        "outcome": "Full editorial estate on EDS. Standard template for all new campaigns.",
    },
    {
        "phase": "Phase 4",
        "label": "M2C Renewal Conversation",
        "time": "Pre-Dec 2027",
        "color": RED,
        "scope": "Use EDS success on editorial layer to frame AEM Cloud Service renewal. EDS = performance tier; AEM CS = complex authoring tier.",
        "outcome": "$9M GNARR potential realized. Expanded AEM CS + EDS deal at renewal.",
    },
]

pw = Inches(2.45)
ph = Inches(4.65)
px = Inches(0.22)
py = Inches(1.2)

for i, p in enumerate(phases):
    cx = px + i * (pw + Inches(0.1))
    # card
    add_rect(slide, cx, py, pw, ph, fill=WHITE, line=p["color"])
    # top color bar
    add_rect(slide, cx, py, pw, Inches(0.6), fill=p["color"])
    # phase label
    add_textbox(slide, p["phase"],
                cx + Inches(0.1), py + Inches(0.05), pw - Inches(0.2), Inches(0.28),
                size=11, bold=True, color=WHITE)
    add_textbox(slide, p["time"],
                cx + Inches(0.1), py + Inches(0.33), pw - Inches(0.2), Inches(0.22),
                size=8, color=RGBColor(0xFF, 0xFF, 0xCC))
    # title
    add_textbox(slide, p["label"],
                cx + Inches(0.12), py + Inches(0.68), pw - Inches(0.2), Inches(0.4),
                size=10, bold=True, color=p["color"], wrap=True)
    # scope
    add_textbox(slide, "Scope:",
                cx + Inches(0.12), py + Inches(1.15), pw - Inches(0.2), Inches(0.22),
                size=8.5, bold=True, color=DARK)
    add_textbox(slide, p["scope"],
                cx + Inches(0.12), py + Inches(1.38), pw - Inches(0.2), Inches(1.55),
                size=8.5, color=MID, wrap=True)
    # divider
    add_rect(slide, cx + Inches(0.12), py + Inches(3.0), pw - Inches(0.25),
             Inches(0.02), fill=RGBColor(0xDD, 0xDD, 0xDD))
    # outcome
    add_textbox(slide, "Outcome:",
                cx + Inches(0.12), py + Inches(3.08), pw - Inches(0.2), Inches(0.22),
                size=8.5, bold=True, color=DARK)
    add_textbox(slide, p["outcome"],
                cx + Inches(0.12), py + Inches(3.3), pw - Inches(0.2), Inches(1.15),
                size=8.5, color=p["color"], bold=False, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — SUCCESS METRICS
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Success Metrics", "Measure before Phase 1 begins; report at Phase 1 completion")

m_headers = ["Metric", "Current Baseline", "EDS Target", "Business Value"]
m_rows = [
    ["Lighthouse Performance Score\n(newsroom, BMH)",
     "Est. 40–60\n(AEM SPA architecture)",
     "100",
     "Better AI indexability; Core Web Vitals compliance; brand halo"],
    ["Time from draft to publish\n(press release)",
     "Days\n(AEM ticket → deployment)",
     "< 30 minutes",
     "Communications team autonomy; faster breaking news response"],
    ["Developer effort per content\nupdate (editorial properties)",
     "Engineering touchpoint\nrequired",
     "Zero\n(author-driven)",
     "Engineering capacity freed for core banking platform"],
    ["Core Web Vitals — LCP\n(newsroom, BMH pages)",
     "Unknown — SPA-rendered\n(likely > 4s)",
     "< 2.5 seconds",
     "User engagement; Google ranking; AI search preference"],
    ["AI search citations for BMH\nfinancial education content",
     "Baseline: LLM Score 80\n(strong vs. peers)",
     "Measurable increase in\nAI platform mentions",
     "Brand authority in AI-generated personal finance answers"],
    ["Content publishing without\ndeveloper involvement",
     "< 20% estimated\n(editorial layer)",
     "> 90%",
     "Reduces IT cost; accelerates publishing cadence"],
]
make_table(slide, m_headers, m_rows,
           Inches(0.3), Inches(1.2),
           Inches(12.7), Inches(5.55),
           col_widths=[2.8, 2.5, 2.4, 4.8],
           header_size=10, body_size=9)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — ACCOUNT STATUS  (INTERNAL ONLY)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Account Status & Financial Considerations",
             "⚠️  INTERNAL USE ONLY — Do not share with the customer")

# Internal warning banner
add_rect(slide, 0, Inches(1.05), W, Inches(0.32), fill=ROW_INT)
add_textbox(slide,
            "⚠️  CONFIDENTIAL — The data on this slide is for internal Adobe use only. "
            "Do not include in any customer-facing presentation.",
            Inches(0.35), Inches(1.08), Inches(12), Inches(0.26),
            size=9, bold=True, color=AMBER)

# Left: financials table
add_textbox(slide, "Account Financials",
            Inches(0.3), Inches(1.52), Inches(6.0), Inches(0.3),
            size=11, bold=True, color=DARK)

fin_headers = ["Field", "Value"]
fin_rows = [
    ["Current Sites ARR (on-prem)",    "$3,200,000"],
    ["Full Adobe Sites ARR",            "~$23M  (ThopseyInternal)"],
    ["Total Adobe Relationship",        "~$29M+  (Sites + Assets + Analytics + AEP + Target + Workfront)"],
    ["Open Pipeline / GNARR Potential","$9,000,000"],
    ["Renewal Date / Quarter",          "December 31, 2027  (FY28 Q1)"],
    ["Products in Use",                 "AEM Sites + Assets (On-Prem), Analytics→CJA, Target, AEP, Doc Cloud"],
    ["Pipeline Wave",                   "In-Progress  (FLM: Herzenberg)"],
    ["M2EDS Readiness",                 "In Progress — \"Neushler & team rebuilt some of their sites on EDS\""],
    ["DR IDs",                          "DR3462836, DR3786061, DR4178248  (+ new: DR4807478, DR4845016, DR4845102)"],
    ["FLM Team",                        "Herzenberg  (FSI portfolio)"],
]
make_table(slide, fin_headers, fin_rows,
           Inches(0.3), Inches(1.88),
           Inches(6.3), Inches(4.5),
           col_widths=[2.5, 3.7],
           header_size=9, body_size=8.5)

# Right: risks
add_textbox(slide, "Risks to the Account",
            Inches(6.85), Inches(1.52), Inches(6.0), Inches(0.3),
            size=11, bold=True, color=DARK)

risks = [
    ("EDS momentum stays informal",
     "Neushler's EDS work risks remaining a side project. Without formal recognition "
     "in the account plan, it won't become the foundation of a renewal success story."),
    ("Active deal cycle creates noise",
     "FLM team managing Target (11 activities), CJA migration, and AEP expansion. "
     "Introduce EDS positioning aligned to the renewal timeline, not the quarterly cycle."),
    ("Competitive CMS evaluation",
     "Unsupported AEM on-prem makes BofA susceptible to Contentful/Optimizely pitches "
     "for the editorial layer. EDS + Document Authoring preempts this conversation."),
]

ry = Inches(1.88)
for title, body in risks:
    add_rect(slide, Inches(6.85), ry, Inches(6.15), Inches(1.35), fill=WHITE,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(slide, Inches(6.85), ry, Inches(0.05), Inches(1.35), fill=AMBER)
    add_textbox(slide, title, Inches(7.05), ry + Inches(0.08),
                Inches(5.9), Inches(0.28), size=10, bold=True, color=DARK)
    add_textbox(slide, body, Inches(7.05), ry + Inches(0.36),
                Inches(5.9), Inches(0.8), size=9, color=MID, wrap=True)
    ry += Inches(1.45)

# Renewal urgency callout
add_rect(slide, Inches(6.85), ry + Inches(0.1), Inches(6.15), Inches(0.95), fill=RED)
add_textbox(slide, "Renewal Urgency",
            Inches(7.0), ry + Inches(0.15), Inches(5.9), Inches(0.25),
            size=9, bold=True, color=WHITE)
add_textbox(slide,
            "December 2027 = 18-month window. Phase 1 (newsroom) live before Dec 2026 "
            "feeds a measured success story into the $9M GNARR renewal.",
            Inches(7.0), ry + Inches(0.42), Inches(5.9), Inches(0.5),
            size=9, color=RGBColor(0xFF, 0xEE, 0xEE), wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — OBJECTION HANDLING  (INTERNAL ONLY)
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Objection Handling",
             "⚠️  INTERNAL USE ONLY — Prepare for these before customer meetings")

add_rect(slide, 0, Inches(1.05), W, Inches(0.32), fill=ROW_INT)
add_textbox(slide,
            "⚠️  CONFIDENTIAL — For internal preparation only. Do not share with the customer.",
            Inches(0.35), Inches(1.08), Inches(12), Inches(0.26),
            size=9, bold=True, color=AMBER)

obj_headers = ["Objection", "Likely Source", "Rebuttal"]
obj_rows = [
    ["\"We invested heavily in AEM on-prem.\nWe can't disrupt the banking platform.\"",
     "Engineering / Platform team",
     "This POV does NOT target the banking portal. EDS is for the editorial layer "
     "(newsroom, BMH, about) only. On-prem stays for the portal. Neushler's team "
     "already proved coexistence works."],
    ["\"Our authors already know AEM tools.\"",
     "AEM power users / digital team",
     "True for the banking portal team. But newsroom and BMH editors are not power "
     "AEM users — they use Word and SharePoint daily. Document Authoring meets them there."],
    ["\"16 DRIDs + complex governance. EDS\ncan't support our compliance needs.\"",
     "Legal / Compliance / IT Security",
     "da.live is built on SharePoint permission models — the same access control and "
     "audit trail infrastructure BofA already uses enterprise-wide. Financial services ready."],
    ["\"The EDS score of 32 means this\nisn't a fit.\"",
     "Internal Adobe stakeholders",
     "The 32 score is for the full site dominated by banking portals. The editorial "
     "layer alone would score 60–70+. The rationale says 'limited marketing-only EDS "
     "scope' — that IS the scope we're targeting."],
    ["\"Neushler's team is already on EDS.\nDo we need a formal engagement?\"",
     "BofA digital / internal stakeholders",
     "Neushler's work is Phase 0 done — great. A formal engagement turns it into an "
     "enterprise platform: official block library, Document Authoring rollout, and a "
     "measured success story for the Dec 2027 renewal."],
    ["\"Performance is already good enough —\nwe have #1 J.D. Power rankings.\"",
     "BofA digital leadership",
     "J.D. Power measures app satisfaction, not editorial Core Web Vitals. newsroom "
     "and BMH are the AI-indexed surfaces. Lighthouse 100 there improves AI search "
     "discoverability — the next digital experience frontier in banking."],
    ["\"Bad timing — we're in CJA migration\nand AEP expansion right now.\"",
     "Account / consulting team concern",
     "EDS on the editorial layer is a parallel track. It doesn't touch Analytics, "
     "AEP, or Target instrumentation. Position as the 'content platform' workstream, "
     "independent of the data program."],
]
make_table(slide, obj_headers, obj_rows,
           Inches(0.3), Inches(1.45),
           Inches(12.7), Inches(5.15),
           col_widths=[2.8, 2.2, 7.6],
           header_size=9, body_size=8.5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — BENEFITS SUMMARY
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Benefits Summary", "What each stakeholder gets from EDS + Document Authoring")

personas = [
    {
        "label": "Business / Marketing",
        "color": RED,
        "bullets": [
            "Press releases live in < 30 minutes from Word — no tickets, no deployments",
            "Better Money Habits articles publish without developer involvement",
            "Campaign microsites launch on EDS — faster time to market",
            "Lighthouse 100 performance = brand credibility + AI search authority",
        ],
    },
    {
        "label": "Authors / Communications",
        "color": BLUE,
        "bullets": [
            "Write in Microsoft Word or Google Docs — no AEM training required",
            "Real-time Sidekick preview before every publish",
            "Bilingual (EN/ES) via SharePoint folders — no separate CMS workflow",
            "Version history and approval flow via SharePoint — nothing new to learn",
        ],
    },
    {
        "label": "Developers / Engineering",
        "color": GREEN,
        "bullets": [
            "Plain HTML, CSS, JavaScript on GitHub — no proprietary SDK",
            "6 standard blocks cover 80% of editorial use cases (built once, shared)",
            "Lighthouse 100 by design — no performance tuning required",
            "17,000 BofA programmers already comfortable with GitHub-native workflow",
        ],
    },
    {
        "label": "IT / Operations",
        "color": AMBER,
        "bullets": [
            "Zero infrastructure to manage for editorial layer (no dispatcher, no app server)",
            "SharePoint-native access controls — compliant with financial services requirements",
            "Static-first delivery reduces security attack surface",
            "Adobe-managed CDN and delivery for editorial properties",
        ],
    },
]

pw = Inches(3.1)
ph = Inches(4.8)
px = Inches(0.25)
py = Inches(1.25)

for i, p in enumerate(personas):
    cx = px + i * (pw + Inches(0.1))
    add_rect(slide, cx, py, pw, ph, fill=WHITE,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(slide, cx, py, pw, Inches(0.07), fill=p["color"])
    add_textbox(slide, p["label"],
                cx + Inches(0.15), py + Inches(0.18), pw - Inches(0.3), Inches(0.36),
                size=12, bold=True, color=p["color"])
    add_rect(slide, cx + Inches(0.15), py + Inches(0.6), pw - Inches(0.3),
             Inches(0.02), fill=RGBColor(0xDD, 0xDD, 0xDD))
    txb = slide.shapes.add_textbox(
        cx + Inches(0.15), py + Inches(0.72), pw - Inches(0.3), ph - Inches(0.9))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for j, b in enumerate(p["bullets"]):
        add_para(tf, f"✓  {b}", size=9.5, color=MID,
                 space_before=(6 if j > 0 else 0))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — RECOMMENDED NEXT STEPS
# ─────────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=LIGHT)
slide_header(slide, "Recommended Next Steps", "Immediate actions for the account team")

steps = [
    {
        "num": "1",
        "title": "Inventory the Neushler EDS Sites",
        "owner": "Account Team + BofA Digital",
        "timing": "Immediately",
        "color": RED,
        "body": ("Get a complete list of sites Neushler's team built on EDS. "
                 "Run Lighthouse. Capture the metrics. "
                 "Declare Phase 0 complete — this becomes the opening slide of every EDS conversation at BofA."),
    },
    {
        "num": "2",
        "title": "Schedule Newsroom Discovery Call",
        "owner": "SC + AE + BofA Communications Lead",
        "timing": "This month",
        "color": BLUE,
        "body": ("Validate the press release publishing pain point with the communications team. "
                 "Run a live Document Authoring demo: draft a press release in Word, "
                 "publish to a demo environment in < 30 minutes. This single demo wins newsroom champions."),
    },
    {
        "num": "3",
        "title": "Better Money Habits Content Audit",
        "owner": "SC",
        "timing": "2–3 weeks",
        "color": GREEN,
        "body": ("Count articles, map URL patterns, identify the interactive tool inventory. "
                 "Confirm the 2 custom blocks are indeed the only scope items beyond standard EDS primitives. "
                 "De-risks Phase 2 and sharpens the SOW."),
    },
    {
        "num": "4",
        "title": "Align with FLM on EDS Timing",
        "owner": "AE + FLM (Herzenberg)",
        "timing": "This quarter",
        "color": AMBER,
        "body": ("Confirm when to introduce EDS alongside the renewal conversation. "
                 "Goal: Phase 1 (newsroom) live before Dec 2026 so it feeds the Dec 2027 "
                 "renewal narrative with 12 months of production data. Coordinate with CJA and AEP timeline."),
    },
    {
        "num": "5",
        "title": "Agree on Baseline Metrics",
        "owner": "SC + BofA Digital",
        "timing": "Before Phase 1 starts",
        "color": SLATE,
        "body": ("Measure now: Lighthouse scores on newsroom and BMH, press release publish time, "
                 "developer touchpoints per editorial update. Without a before snapshot, "
                 "the after story has no contrast. This is the foundation of the renewal case."),
    },
    {
        "num": "6",
        "title": "Identify FSI EDS Reference Customer",
        "owner": "SC + Adobe Marketing",
        "timing": "This quarter",
        "color": MID,
        "body": ("Find a financial services reference customer running EDS for editorial/communications content. "
                 "JPMorgan Chase context (same AEM architecture complexity) would resonate strongly. "
                 "A peer reference removes the 'are we the guinea pig?' objection."),
    },
]

step_w = Inches(4.05)
step_h = Inches(2.3)
col_gap = Inches(0.2)
row_gap = Inches(0.18)

for i, s in enumerate(steps):
    col = i % 3
    row = i // 3
    cx = Inches(0.22) + col * (step_w + col_gap)
    cy = Inches(1.25) + row * (step_h + row_gap)

    add_rect(slide, cx, cy, step_w, step_h, fill=WHITE,
             line=RGBColor(0xDD, 0xDD, 0xDD))
    # left color bar
    add_rect(slide, cx, cy, Inches(0.06), step_h, fill=s["color"])
    # number
    add_rect(slide, cx + Inches(0.18), cy + Inches(0.1),
             Inches(0.44), Inches(0.44), fill=s["color"])
    add_textbox(slide, s["num"],
                cx + Inches(0.18), cy + Inches(0.1), Inches(0.44), Inches(0.44),
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # title
    add_textbox(slide, s["title"],
                cx + Inches(0.76), cy + Inches(0.14), step_w - Inches(0.9), Inches(0.4),
                size=10.5, bold=True, color=DARK, wrap=True)
    # owner + timing
    add_textbox(slide,
                f"Owner: {s['owner']}  |  {s['timing']}",
                cx + Inches(0.76), cy + Inches(0.52), step_w - Inches(0.9), Inches(0.22),
                size=7.5, italic=True, color=s["color"])
    # body
    add_textbox(slide, s["body"],
                cx + Inches(0.18), cy + Inches(0.82), step_w - Inches(0.3), step_h - Inches(0.92),
                size=8.5, color=MID, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────────────────
out_path = "/Users/jihuang/Desktop/skills/plugins/aem/edge-delivery-services/skills/customer-pov/output/bank-of-america-pov-2026-05-27.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
